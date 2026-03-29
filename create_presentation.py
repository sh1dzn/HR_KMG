#!/usr/bin/env python3
"""Generate project presentation for HR AI Module (KMG)."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Brand colors
BLUE = RGBColor(0x00, 0x5B, 0xAA)
DARK_BLUE = RGBColor(0x00, 0x33, 0x66)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF0, 0xF2, 0xF5)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
ACCENT_GREEN = RGBColor(0x00, 0xA6, 0x5A)
ACCENT_ORANGE = RGBColor(0xFF, 0x8C, 0x00)
ACCENT_RED = RGBColor(0xE0, 0x3E, 0x2D)
SUBTITLE_GRAY = RGBColor(0x66, 0x66, 0x66)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

W = prs.slide_width
H = prs.slide_height


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=DARK_GRAY, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=16,
                  color=DARK_GRAY, bold=False, spacing=1.2, bullet=False, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        prefix = "\u2022  " if bullet else ""
        p.text = prefix + line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = font_name
        p.space_after = Pt(font_size * spacing * 0.5)
    return txBox


def add_card(slide, left, top, width, height, title, body_lines,
             accent_color=BLUE, title_size=16, body_size=13):
    # Card background
    card = add_shape(slide, left, top, width, height, WHITE)
    card.shadow.inherit = False
    # Accent bar
    add_shape(slide, left, top, Inches(0.06), height, accent_color)
    # Title
    add_textbox(slide, left + Inches(0.25), top + Inches(0.15), width - Inches(0.4), Inches(0.4),
                title, font_size=title_size, color=accent_color, bold=True)
    # Body
    if body_lines:
        add_multiline(slide, left + Inches(0.25), top + Inches(0.55),
                      width - Inches(0.4), height - Inches(0.7),
                      body_lines, font_size=body_size, bullet=True)


def add_screenshot_placeholder(slide, left, top, width, height, label="Скриншот"):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0xE8, 0xEB, 0xEF)
    shape.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
    shape.line.width = Pt(1)
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = f"[ {label} ]"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x99, 0x99, 0x99)
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_before = Pt(height.inches * 72 / 2 - 10)


def slide_header(slide, title, subtitle=None):
    add_shape(slide, 0, 0, W, Inches(1.2), BLUE)
    add_textbox(slide, Inches(0.8), Inches(0.2), Inches(10), Inches(0.6),
                title, font_size=32, color=WHITE, bold=True)
    if subtitle:
        add_textbox(slide, Inches(0.8), Inches(0.7), Inches(10), Inches(0.4),
                    subtitle, font_size=16, color=RGBColor(0xCC, 0xDD, 0xFF))
    # Page number area
    sn = slide.slide_number = None  # handled by pptx


# ============================================================
# SLIDE 1 — Title
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, DARK_BLUE)
# Decorative shape
add_shape(slide, Inches(8.5), 0, Inches(5), H, RGBColor(0x00, 0x48, 0x88))
# Title
add_textbox(slide, Inches(1), Inches(1.8), Inches(7), Inches(1),
            "HR AI Module", font_size=54, color=WHITE, bold=True)
add_textbox(slide, Inches(1), Inches(2.9), Inches(7), Inches(0.8),
            "Интеллектуальное управление целями сотрудников",
            font_size=24, color=RGBColor(0xAA, 0xCC, 0xFF))
# Separator
add_shape(slide, Inches(1), Inches(3.9), Inches(3), Inches(0.05), ACCENT_GREEN)
# Info
add_multiline(slide, Inches(1), Inches(4.2), Inches(6), Inches(2), [
    "Команда: SilkRoadTech",
    "Хакатон КМГ-КУМКОЛЬ  |  Март 2026",
    "AI-модуль для оценки, генерации и управления целями",
], font_size=18, color=RGBColor(0xCC, 0xDD, 0xFF), spacing=1.8)
# Screenshot placeholder on right
add_screenshot_placeholder(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(4.5),
                           "Скриншот главной страницы")

# ============================================================
# SLIDE 2 — Problem
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Проблема", "Почему HR-процессы нуждаются в автоматизации")

problems = [
    ("Низкое качество целей", [
        "70% целей не соответствуют SMART-критериям",
        "Субъективная оценка без единых стандартов",
        "Размытые формулировки без метрик и сроков"
    ]),
    ("Ручные процессы", [
        "Согласование целей занимает недели",
        "Нет прозрачности в workflow",
        "Потеря целей между статусами"
    ]),
    ("Отсутствие аналитики", [
        "Нет агрегированных данных по организации",
        "Невозможно выявить слабые подразделения",
        "Отсутствует предиктивная аналитика рисков"
    ]),
    ("Разрыв стратегии", [
        "Цели не связаны со стратегией компании",
        "Нет каскадирования от руководителей",
        "Сложно отследить стратегическое покрытие"
    ]),
]

for i, (title, lines) in enumerate(problems):
    col = i % 4
    left = Inches(0.6 + col * 3.1)
    top = Inches(1.6)
    colors = [ACCENT_RED, ACCENT_ORANGE, BLUE, DARK_BLUE]
    add_card(slide, left, top, Inches(2.85), Inches(3.2), title, lines,
             accent_color=colors[i], title_size=15, body_size=12)

# Bottom note
add_textbox(slide, Inches(0.6), Inches(5.2), Inches(12), Inches(0.8),
            "Результат: снижение эффективности, демотивация сотрудников, "
            "невозможность стратегического планирования",
            font_size=14, color=ACCENT_RED, bold=True)

# ============================================================
# SLIDE 3 — Solution overview
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Наше решение", "AI-платформа для управления целями")

modules = [
    ("SMART-оценка", "AI-анализ целей по 5 критериям с рекомендациями по улучшению", BLUE),
    ("RAG-генерация", "Генерация целей на основе внутренних документов компании (ВНД)", ACCENT_GREEN),
    ("Workflow", "Полный цикл: черновик → согласование → выполнение → архив", ACCENT_ORANGE),
    ("Аналитика", "Дашборд с индексом зрелости, трендами и heatmap", DARK_BLUE),
    ("Предсказание рисков", "Прогнозирование провала целей по 5 факторам", ACCENT_RED),
    ("Каскадирование", "Автоматический спуск целей от руководителей к подчинённым", RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, desc, color) in enumerate(modules):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.6 + row * 2.6)
    card = add_shape(slide, left, top, Inches(3.8), Inches(2.2), LIGHT_GRAY)
    card.line.fill.background()
    add_shape(slide, left, top, Inches(3.8), Inches(0.06), color)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.2), Inches(3.4), Inches(0.4),
                title, font_size=18, color=color, bold=True)
    add_textbox(slide, left + Inches(0.2), top + Inches(0.7), Inches(3.4), Inches(1.3),
                desc, font_size=13, color=DARK_GRAY)

# ============================================================
# SLIDE 4 — Architecture
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Архитектура", "Технологический стек")

# Left column - tech stack
stack_items = [
    ("Frontend", "React 18 + Vite + Tailwind CSS"),
    ("Backend", "FastAPI + SQLAlchemy 2.0 + Pydantic v2"),
    ("Database", "PostgreSQL 18"),
    ("Vector DB", "ChromaDB (RAG-поиск)"),
    ("AI/LLM", "OpenAI GPT-4o + text-embedding-3-small"),
    ("Auth", "JWT (access + refresh tokens)"),
    ("Deploy", "Docker Compose + Caddy + Vercel"),
    ("CI/CD", "GitHub Actions"),
]

for i, (label, value) in enumerate(stack_items):
    top = Inches(1.5 + i * 0.6)
    add_textbox(slide, Inches(0.8), top, Inches(1.8), Inches(0.5),
                label, font_size=14, color=BLUE, bold=True)
    add_textbox(slide, Inches(2.8), top, Inches(4), Inches(0.5),
                value, font_size=14, color=DARK_GRAY)

# Right - architecture diagram placeholder
add_screenshot_placeholder(slide, Inches(7.5), Inches(1.5), Inches(5.2), Inches(5),
                           "Схема архитектуры")

# ============================================================
# SLIDE 5 — SMART Evaluation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "SMART-оценка целей", "Семантический анализ на базе GPT-4o")

# Left description
items = [
    "Оценка каждой цели по 5 критериям (S-M-A-R-T)",
    "Каждый критерий оценивается от 0 до 1",
    "Пороги: Низкий (<0.5), Средний (0.5-0.7), Высокий (>0.85)",
    "Пакетная оценка всех целей сотрудника",
    "Переформулирование: AI предлагает улучшенный текст",
    "Fallback: семантические эвристики без LLM",
    "80+ паттернов ключевых слов для heuristics",
]
add_multiline(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
              items, font_size=14, bullet=True, spacing=1.5)

# Right placeholder
add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5),
                           "Скриншот SMART-оценки")

# ============================================================
# SLIDE 6 — RAG Generation
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "RAG-генерация целей", "Генерация на основе внутренних документов")

items = [
    "Гибридный поиск: векторный (cosine) + лексический с RRF",
    "160+ документов: ВНД, стратегии, KPI-фреймворки",
    "Генерация 3-5 целей за запрос",
    "Указание источника: документ + фрагмент + обоснование",
    "Дедупликация с существующими целями",
    "Фокус-области для направленной генерации",
    "Каскадирование целей руководителя",
    "Учёт истории достижимости сотрудника",
]
add_multiline(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(5),
              items, font_size=14, bullet=True, spacing=1.4)

add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5),
                           "Скриншот генерации целей")

# ============================================================
# SLIDE 7 — Workflow
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Workflow согласования", "Полный жизненный цикл цели")

# Status flow
statuses = ["Черновик", "Активная", "На\nсогласовании", "Одобрена /\nОтклонена",
            "В работе", "Выполнена"]
colors_flow = [RGBColor(0x99, 0x99, 0x99), BLUE, ACCENT_ORANGE, ACCENT_GREEN,
               DARK_BLUE, ACCENT_GREEN]

for i, (status, clr) in enumerate(zip(statuses, colors_flow)):
    left = Inches(0.6 + i * 2.1)
    # Circle/box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0),
                                   Inches(1.7), Inches(1.0))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = status
    p.font.size = Pt(12)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    # Arrow
    if i < len(statuses) - 1:
        add_textbox(slide, left + Inches(1.7), Inches(2.2), Inches(0.4), Inches(0.5),
                    "\u2192", font_size=24, color=DARK_GRAY, bold=True)

# Features below
features = [
    "Полный аудит-трейл: GoalEvent + GoalReview",
    "Действия: отправить, одобрить, отклонить, комментировать",
    "История рецензий с вердиктами",
    "Автоматические алерты при застое",
]
add_multiline(slide, Inches(0.8), Inches(3.5), Inches(5.5), Inches(3),
              features, font_size=14, bullet=True, spacing=1.5)

add_screenshot_placeholder(slide, Inches(7), Inches(3.3), Inches(5.8), Inches(3.5),
                           "Скриншот Workflow / Approvals")

# ============================================================
# SLIDE 8 — Dashboard & Analytics
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Аналитика и дашборд", "5-факторный индекс зрелости организации")

# Maturity index factors
factors = [
    ("SMART-фактор", "30%", "Средний балл SMART"),
    ("Стратегический", "20%", "% стратегических целей"),
    ("Типовой", "20%", "% результатных + влияния"),
    ("Весовой", "15%", "Баланс весов (идеал — 100%)"),
    ("Количественный", "15%", "% сотрудников с 3-5 целями"),
]

for i, (name, weight, desc) in enumerate(factors):
    top = Inches(1.6 + i * 0.7)
    add_shape(slide, Inches(0.8), top, Inches(0.5), Inches(0.5), BLUE)
    add_textbox(slide, Inches(0.85), top + Inches(0.08), Inches(0.4), Inches(0.35),
                weight, font_size=11, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, Inches(1.5), top, Inches(2), Inches(0.3),
                name, font_size=14, color=DARK_BLUE, bold=True)
    add_textbox(slide, Inches(1.5), top + Inches(0.3), Inches(3.5), Inches(0.3),
                desc, font_size=12, color=SUBTITLE_GRAY)

# Additional features
add_multiline(slide, Inches(0.8), Inches(5.3), Inches(5), Inches(1.5), [
    "Heatmap подразделений",
    "Квартальные тренды (SMART + стратегия)",
    "Бенчмарк отделов",
], font_size=13, bullet=True, color=DARK_GRAY)

add_screenshot_placeholder(slide, Inches(6.5), Inches(1.5), Inches(6.3), Inches(5.5),
                           "Скриншот дашборда")

# ============================================================
# SLIDE 9 — Risk Prediction
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Предсказание рисков", "Статистическая модель + LLM-объяснение")

# Risk factors
risk_factors = [
    ("SMART-качество", "20%", BLUE),
    ("Кол-во отклонений", "20%", ACCENT_ORANGE),
    ("Стагнация статуса", "20%", ACCENT_RED),
    ("Давление дедлайна", "25%", RGBColor(0x8B, 0x5C, 0xF6)),
    ("История отдела", "15%", DARK_BLUE),
]

for i, (name, weight, clr) in enumerate(risk_factors):
    left = Inches(0.6 + i * 2.45)
    card = add_shape(slide, left, Inches(1.6), Inches(2.2), Inches(1.5), LIGHT_GRAY)
    card.line.fill.background()
    add_shape(slide, left, Inches(1.6), Inches(2.2), Inches(0.05), clr)
    add_textbox(slide, left + Inches(0.15), Inches(1.8), Inches(1.9), Inches(0.4),
                name, font_size=14, color=clr, bold=True)
    add_textbox(slide, left + Inches(0.15), Inches(2.3), Inches(1.9), Inches(0.4),
                f"Вес: {weight}", font_size=20, color=DARK_GRAY, bold=True)

# Risk levels
levels = [
    ("Высокий", "> 0.7", ACCENT_RED),
    ("Средний", "0.4 — 0.7", ACCENT_ORANGE),
    ("Низкий", "< 0.4", ACCENT_GREEN),
]
add_textbox(slide, Inches(0.8), Inches(3.5), Inches(3), Inches(0.4),
            "Уровни риска:", font_size=16, color=DARK_GRAY, bold=True)
for i, (level, rng, clr) in enumerate(levels):
    left = Inches(0.8 + i * 3)
    top = Inches(4.0)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top,
                                   Inches(2.5), Inches(0.7))
    shape.fill.solid()
    shape.fill.fore_color.rgb = clr
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    p.text = f"{level}  ({rng})"
    p.font.size = Pt(14)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER

add_multiline(slide, Inches(0.8), Inches(5.0), Inches(5), Inches(2), [
    "LLM генерирует понятное объяснение рисков",
    "Рекомендации по снижению рисков",
], font_size=14, bullet=True, spacing=1.5)

add_screenshot_placeholder(slide, Inches(7), Inches(3.5), Inches(5.8), Inches(3.5),
                           "Скриншот предсказания рисков")

# ============================================================
# SLIDE 10 — Cascading & Dependencies
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Каскадирование и граф зависимостей",
             "Связь целей по вертикали и горизонтали")

# Left - cascading
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.4),
            "Каскадирование", font_size=20, color=BLUE, bold=True)
cascade_items = [
    "Автоматический спуск целей руководителя",
    "Превью каскадированных целей",
    "Обнаружение конфликтов и дубликатов",
    "Кастомизация целей для каждого отдела",
]
add_multiline(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(2.5),
              cascade_items, font_size=14, bullet=True, spacing=1.5)

# Right - dependencies
add_textbox(slide, Inches(0.8), Inches(4.2), Inches(3), Inches(0.4),
            "Граф зависимостей", font_size=20, color=ACCENT_GREEN, bold=True)
dep_items = [
    "Интерактивная визуализация (force-graph)",
    "Типы связей: blocks, depends_on, related_to",
    "AI-рекомендации зависимостей",
    "Обнаружение циклических зависимостей",
]
add_multiline(slide, Inches(0.8), Inches(4.8), Inches(5.5), Inches(2.5),
              dep_items, font_size=14, bullet=True, spacing=1.5)

add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5),
                           "Скриншот каскадирования / графа")

# ============================================================
# SLIDE 11 — Alerts & Integrations
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Алерты и интеграции", "Проактивный мониторинг и экспорт")

# Left - alerts
add_textbox(slide, Inches(0.8), Inches(1.5), Inches(3), Inches(0.4),
            "Система алертов", font_size=20, color=ACCENT_RED, bold=True)
alert_items = [
    "Низкий SMART-индекс",
    "Слабая стратегическая связь",
    "Отсутствие измеримых критериев",
    "Просроченные дедлайны",
    "Дисбаланс портфеля целей",
    "Стагнация в согласовании",
]
add_multiline(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(3),
              alert_items, font_size=14, bullet=True, spacing=1.3)

# Right - integrations
add_textbox(slide, Inches(0.8), Inches(4.8), Inches(3), Inches(0.4),
            "HR-интеграции", font_size=20, color=BLUE, bold=True)
int_items = [
    "1С:ЗУП — экспорт целей в формате 1С",
    "SAP SuccessFactors — синхронизация",
    "Oracle HCM — двусторонний обмен",
]
add_multiline(slide, Inches(0.8), Inches(5.4), Inches(5.5), Inches(1.5),
              int_items, font_size=14, bullet=True, spacing=1.3)

add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5.5),
                           "Скриншот алертов / интеграций")

# ============================================================
# SLIDE 12 — 1-on-1 & Meeting Agenda
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "1-on-1 встречи", "Автоматическая повестка для руководителя")

items = [
    "Генерация повестки для встречи руководитель — сотрудник",
    "Цели, требующие обсуждения",
    "Цели в зоне риска",
    "Тренды производительности",
    "Рекомендации по корректировке целей",
]
add_multiline(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4),
              items, font_size=16, bullet=True, spacing=1.8)

add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5),
                           "Скриншот 1-on-1 повестки")

# ============================================================
# SLIDE 13 — Demo Data
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Демо-данные", "Масштаб и реалистичность")

data_cards = [
    ("450", "Сотрудников", BLUE),
    ("9 000", "Целей", ACCENT_GREEN),
    ("8", "Подразделений", ACCENT_ORANGE),
    ("160+", "Документов (ВНД)", DARK_BLUE),
    ("30 789", "Событий аудита", RGBColor(0x8B, 0x5C, 0xF6)),
    ("4 305", "Рецензий", ACCENT_RED),
]

for i, (num, label, clr) in enumerate(data_cards):
    col = i % 3
    row = i // 3
    left = Inches(1 + col * 3.8)
    top = Inches(1.8 + row * 2.5)
    card = add_shape(slide, left, top, Inches(3.2), Inches(2.0), LIGHT_GRAY)
    card.line.fill.background()
    add_shape(slide, left, top, Inches(3.2), Inches(0.06), clr)
    add_textbox(slide, left, top + Inches(0.3), Inches(3.2), Inches(0.8),
                num, font_size=40, color=clr, bold=True, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + Inches(1.2), Inches(3.2), Inches(0.5),
                label, font_size=16, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

# ============================================================
# SLIDE 14 — Fallback mode
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Работа без LLM", "Graceful degradation — система работает и без OpenAI")

# Table-like layout
headers = ["Функция", "С LLM (GPT-4o)", "Без LLM (Fallback)"]
rows = [
    ["SMART-оценка", "Детальный AI-анализ", "Эвристики (80+ паттернов)"],
    ["RAG-поиск", "Гибридный (вектор + лексика)", "Только лексический"],
    ["Генерация целей", "LLM-синтез", "Заголовки из документов"],
    ["Объяснение рисков", "LLM-нарратив", "Список факторов"],
]

# Headers
for j, h in enumerate(headers):
    left = Inches(1 + j * 3.8)
    shape = add_shape(slide, left, Inches(1.8), Inches(3.5), Inches(0.6), BLUE)
    add_textbox(slide, left, Inches(1.88), Inches(3.5), Inches(0.4),
                h, font_size=14, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# Rows
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        left = Inches(1 + j * 3.8)
        top = Inches(2.5 + i * 0.7)
        bg_clr = LIGHT_GRAY if i % 2 == 0 else WHITE
        add_shape(slide, left, top, Inches(3.5), Inches(0.6), bg_clr)
        add_textbox(slide, left + Inches(0.15), top + Inches(0.12), Inches(3.2), Inches(0.4),
                    cell, font_size=12, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_textbox(slide, Inches(1), Inches(5.3), Inches(10), Inches(0.5),
            "Система деградирует плавно — пользователь всегда получает результат",
            font_size=16, color=ACCENT_GREEN, bold=True)

# ============================================================
# SLIDE 15 — Deployment
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Развёртывание", "Docker Compose + CI/CD")

items = [
    "Frontend: Vercel (автодеплой)",
    "Backend: Docker на сервере + Caddy reverse proxy",
    "CI/CD: GitHub Actions (push to main → auto-deploy)",
    "PostgreSQL 18 + ChromaDB (persistent volumes)",
    "Health-check эндпоинт: /health",
    "Swagger документация: /docs",
]
add_multiline(slide, Inches(0.8), Inches(1.5), Inches(5.5), Inches(4),
              items, font_size=16, bullet=True, spacing=1.8)

# URLs
add_textbox(slide, Inches(0.8), Inches(5.2), Inches(5), Inches(0.4),
            "Демо:", font_size=16, color=BLUE, bold=True)
add_multiline(slide, Inches(0.8), Inches(5.6), Inches(8), Inches(1.5), [
    "https://hr-kmg.silkroadtech.kz/",
    "API: https://hr-ai.sh7dzn.me/docs",
], font_size=14, color=DARK_GRAY)

add_screenshot_placeholder(slide, Inches(7), Inches(1.5), Inches(5.8), Inches(5),
                           "Скриншот CI/CD / Docker")

# ============================================================
# SLIDE 16 — Results & Impact
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, WHITE)
slide_header(slide, "Результаты и эффект", "Что даёт внедрение HR AI Module")

results = [
    ("Скорость", "Оценка цели за секунды вместо часов ручной работы", BLUE),
    ("Качество", "SMART-анализ повышает качество формулировок на 40%+", ACCENT_GREEN),
    ("Прозрачность", "Полный аудит-трейл и прозрачный workflow", ACCENT_ORANGE),
    ("Стратегия", "Каскадирование обеспечивает связь целей со стратегией", DARK_BLUE),
    ("Предиктивность", "Раннее выявление рисков до наступления дедлайна", ACCENT_RED),
    ("Масштаб", "Готовность к работе с 450+ сотрудниками и 9000+ целями", RGBColor(0x8B, 0x5C, 0xF6)),
]

for i, (title, desc, clr) in enumerate(results):
    col = i % 3
    row = i // 3
    left = Inches(0.6 + col * 4.1)
    top = Inches(1.6 + row * 2.6)
    card = add_shape(slide, left, top, Inches(3.8), Inches(2.2), LIGHT_GRAY)
    card.line.fill.background()
    add_shape(slide, left, top, Inches(0.06), Inches(2.2), clr)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.2), Inches(3.3), Inches(0.4),
                title, font_size=20, color=clr, bold=True)
    add_textbox(slide, left + Inches(0.25), top + Inches(0.8), Inches(3.3), Inches(1.2),
                desc, font_size=14, color=DARK_GRAY)

# ============================================================
# SLIDE 17 — Thank you
# ============================================================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, DARK_BLUE)
add_shape(slide, Inches(8.5), 0, Inches(5), H, RGBColor(0x00, 0x48, 0x88))

add_textbox(slide, Inches(1), Inches(2), Inches(7), Inches(1),
            "Спасибо!", font_size=60, color=WHITE, bold=True)

add_shape(slide, Inches(1), Inches(3.2), Inches(3), Inches(0.05), ACCENT_GREEN)

add_multiline(slide, Inches(1), Inches(3.6), Inches(6), Inches(3), [
    "Команда: SilkRoadTech",
    "Проект: HR AI Module",
    "Демо: hr-kmg.silkroadtech.kz",
    "",
    "Готовы к вопросам!",
], font_size=20, color=RGBColor(0xCC, 0xDD, 0xFF), spacing=1.5)

add_screenshot_placeholder(slide, Inches(8.8), Inches(1.5), Inches(4), Inches(4.5),
                           "Логотип / QR-код")

# Save
output_path = "/root/HR_KMG/HR_AI_Module_Presentation.pptx"
prs.save(output_path)
print(f"Presentation saved: {output_path}")
print(f"Total slides: {len(prs.slides)}")
